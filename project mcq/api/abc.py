from pptx import Presentation
import glob
emp_list=[]
for eachfile in glob.glob("Rigil Capabilities External  (1)-converted.pptx"):
    prs = Presentation(eachfile)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                emp_list.append(shape.text)

def text_cleaning(text_data):
    import re
    import string
    text_data=re.sub(r'\s+', ' ', text_data)
    text_data=re.sub(r'[0-9]+', '', text_data)
    text_data=re.sub(r',','',text_data)
    text_data=re.sub('<[^>]+?>', '', text_data)
    text_data = "".join([ch for ch in text_data if ch.isalnum() or ch in string.punctuation or ch.isspace()])
    text_data=text_data.replace('www.','').replace('.com','').replace('rigil','').replace('&','and').replace(' Lets make life better!','').replace(":",'')
    text_data = text_data.encode("ascii", "ignore")
    text_data = text_data.decode()
    return text_data
list_emp=[]
for i in emp_list[0:-22]:
    list_emp.append(text_cleaning(i))
for j in list_emp:
    if j=="" or j==" " or j=="  ":
        list_emp.remove(j)
text_data_emp=[]
def ph_to_se():
    from keytotext import pipeline
    nlp = pipeline("mrm8488/t5-base-finetuned-common_gen")  #use k2t(231mb) or k2t-base(1gb) or k2t-tiny or mrm8488/t5-base-finetuned-common_gen(1.11 GB)
    params = {"do_sample":True, "num_beams":3, "no_repeat_ngram_size":2, "early_stopping":True}  
    for i in list_emp: 
        #print(list_emp[1])
        text_data_emp.append(nlp([i.split(" ")]),**params)
ph_to_se()

emp_data=" ".join(text_data_emp)
print(emp_data)
    







 


                